"""Generate HFpEF findings presentation as .pptx."""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent

# Colors
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
MID_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
TABLE_HEADER_BG = RGBColor(0x00, 0x56, 0x8A)
TABLE_ROW_LIGHT = RGBColor(0xF5, 0xF9, 0xFC)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN_ACCENT = RGBColor(0x2E, 0xCC, 0x71)
ORANGE_ACCENT = RGBColor(0xF3, 0x9C, 0x12)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=13,
                      color=DARK_TEXT, font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, sz, clr) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz if sz else font_size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(4)
    return txBox


def add_table(slide, left, top, width, height, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)

            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT

            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_WHITE

    return table_shape


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # =========================================================================
    # SLIDE 1: Title
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.2),
                 "HFpEF Protein-Disease Association Classification",
                 font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
                 "Automating Literature Curation with Biomedical NLP",
                 font_size=22, color=ACCENT_TEAL, bold=False, alignment=PP_ALIGN.LEFT)

    # Thin accent line
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(3.0), Inches(3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.6),
                 "Aktan Azat  |  February 2026",
                 font_size=16, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

    # =========================================================================
    # SLIDE 2: The Problem
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "The Problem", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    lines = [
        ("CaseOLAP identifies 3,554 candidate proteins from PubMed literature, ranked by text co-occurrence with HFpEF.", False, 15, DARK_TEXT),
        ("", False, 8, DARK_TEXT),
        ("The fundamental problem: co-occurrence is not evidence.", True, 16, RED_ACCENT),
        ("", False, 8, DARK_TEXT),
        ("A protein gets a high CaseOLAP score simply by appearing in the same abstract as HFpEF -- but the actual sentence might say:", False, 14, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("  \"BNP levels were elevated in HFpEF patients\"              -->  True association", False, 13, GREEN_ACCENT),
        ("  \"Patients with HFpEF were excluded from this study\"     -->  Incidental mention", False, 13, ORANGE_ACCENT),
        ("  \"No association was found between TNF-alpha and HFpEF\" -->  Negative finding", False, 13, RED_ACCENT),
        ("", False, 8, DARK_TEXT),
        ("Manual review of thousands of protein-paper pairs is not feasible.", False, 14, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Goal: Build an NLP classifier to automatically determine whether a PubMed sentence describes a real association, a negative finding, or an incidental mention.", True, 15, ACCENT_BLUE),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.5), lines)

    # =========================================================================
    # SLIDE 3: The Pipeline
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "The Pipeline", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    steps = [
        ("1", "MeSH Terms + CaseOLAP Proteins", "7 HFpEF MeSH descriptors + 3,554 protein IDs with UniProt synonyms"),
        ("2", "PubMed Search", "NCBI E-utilities API: disease terms AND protein names; ~25,000 abstracts retrieved"),
        ("3", "Sentence Extraction", "Identify sentences with protein-disease co-mentions; section header stripping"),
        ("4", "NLP Classification", "PubMedBERT classifier: associated | not_associated | incidental"),
        ("5", "Filtered Rankings", "Evidence-based protein rankings with per-protein confidence scores"),
    ]

    y = Inches(1.5)
    for num, title, desc in steps:
        # Number circle
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(1.2), y, Inches(0.55), Inches(0.55))
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_BLUE
        shape.line.fill.background()
        tf = shape.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(2.0), y - Inches(0.05), Inches(4), Inches(0.4),
                     title, font_size=16, color=DARK_TEXT, bold=True)
        add_text_box(slide, Inches(2.0), y + Inches(0.35), Inches(9), Inches(0.4),
                     desc, font_size=12, color=MID_GRAY)

        # Arrow
        if num != "5":
            arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                           Inches(1.33), y + Inches(0.6), Inches(0.3), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_GRAY
            arrow.line.fill.background()

        y += Inches(1.1)

    # =========================================================================
    # SLIDE 4: The Dataset
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "The Dataset", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    lines = [
        ("1,168 sentences manually labeled from PubMed abstracts co-mentioning a protein and HFpEF", False, 15, DARK_TEXT),
        ("Split: 992 training / 176 held-out evaluation (stratified, never seen during training)", False, 13, MID_GRAY),
        ("", False, 6, DARK_TEXT),
        ("Multiple labeling rounds with balanced sampling and negative enrichment", False, 14, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(1.2), Inches(5.5), Inches(2), lines)

    # Dataset composition table
    rows = [
        ["Label", "Definition", "Train", "Eval", "% of Data"],
        ["Associated", "Real protein-disease relationship described", "322", "57", "32.5%"],
        ["Not Associated", "Explicitly states no relationship found", "89", "16", "9.0%"],
        ["Incidental", "Protein mentioned, no claim made", "581", "103", "58.6%"],
    ]
    add_table(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(2.2), rows,
              col_widths=[Inches(1.8), Inches(4.5), Inches(1.2), Inches(1.2), Inches(1.2)])

    lines2 = [
        ("Key challenge: severe class imbalance", True, 15, RED_ACCENT),
        ("\"Not associated\" (explicit negative findings) = only 9% of data -- these are underreported in literature", False, 13, DARK_TEXT),
        ("Addressed with: data augmentation (89 -> 270 samples) + focal loss function", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(1.5), lines2)

    # =========================================================================
    # SLIDE 5: The Models
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "The Models", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    model_rows = [
        ["Model", "Source", "Pre-training Data", "Parameters", "Key Characteristic"],
        ["PubMedBERT", "Microsoft", "PubMed abstracts + full-text", "110M", "Domain-specific vocabulary; trained from scratch on biomedical text"],
        ["SciBERT", "Allen AI", "1.14M scientific papers", "110M", "Broader scientific coverage; custom scivocab"],
        ["BiomedBERT", "Microsoft", "PubMed abstracts only", "110M / 340M", "PubMedBERT successor; large variant available"],
    ]
    add_table(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5), model_rows,
              col_widths=[Inches(1.6), Inches(1.2), Inches(2.8), Inches(1.3), Inches(4.6)])

    lines = [
        ("Training Configuration", True, 16, ACCENT_BLUE),
        ("", False, 4, DARK_TEXT),
        ("Loss: Focal Loss -- FL(p_t) = -alpha_t * (1 - p_t)^gamma * log(p_t)", False, 13, DARK_TEXT),
        ("  Down-weights easy examples; focuses learning on hard, misclassified cases (gamma=2.0)", False, 12, MID_GRAY),
        ("Class weights: Inverse frequency -- incidental: 0.67, associated: 1.21, not_associated: 1.45", False, 13, DARK_TEXT),
        ("Learning rate: Swept {1e-5, 2e-5, 3e-5}; best = 1e-5 (more stable for small datasets)", False, 13, DARK_TEXT),
        ("Early stopping: Patience of 3 epochs on validation F1; up to 8 epochs max", False, 13, DARK_TEXT),
        ("Ensemble: Softmax probability averaging across top 3 models", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.2), lines)

    # =========================================================================
    # SLIDE 6: Excluded Proteins
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Excluded Proteins from CaseOLAP", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "500 proteins (14%) filtered out -- zero associated mentions across all evidence sentences",
                 font_size=15, color=RED_ACCENT, bold=True)

    # Columns: Protein | UniProt | CaseOLAP Score | Assoc | Not Assoc | Incidental | Model Conf
    excl_rows = [
        ["Protein", "UniProt ID", "CaseOLAP", "Assoc", "Not Assoc", "Incidental", "Avg Conf"],
        ["SPARC", "P09486", "0.847", "0", "0", "0", "--"],
        ["Stanniocalcin-1", "P52823", "0.847", "0", "0", "0", "--"],
        ["BNP", "P16860", "--", "0", "0", "12", "0.61"],
        ["TNF-alpha", "P01375", "--", "0", "1", "8", "0.54"],
        ["IL-6", "P05231", "--", "0", "0", "6", "0.58"],
        ["CRP", "P02741", "--", "0", "0", "5", "0.55"],
        ["ACE", "P12821", "0.494", "0", "0", "4", "0.52"],
        ["VEGF-A", "P15692", "--", "0", "0", "3", "0.53"],
        ["Endothelin-1", "P05305", "--", "0", "0", "0", "--"],
        ["Insulin", "P01308", "0.495", "0", "0", "7", "0.56"],
        ["Troponin T", "P45379", "--", "0", "0", "9", "0.59"],
        ["TLR3", "O15455", "0.611", "0", "0", "0", "--"],
        ["Thrombopoietin", "P40225", "0.611", "0", "0", "0", "--"],
        ["Albumin", "P02768", "0.566", "0", "0", "4", "0.54"],
        ["Renin", "P00797", "0.497", "0", "0", "3", "0.51"],
    ]
    add_table(slide, Inches(0.4), Inches(1.8), Inches(12.5), Inches(5.2), excl_rows,
              col_widths=[Inches(2.0), Inches(1.3), Inches(1.2), Inches(1.0), Inches(1.3), Inches(1.3), Inches(1.3)])

    lines_excl = [
        ("Exclusion criterion: associated_mentions == 0 across all evidence sentences", False, 13, MID_GRAY),
        ("\"--\" = no protein-disease co-mention found in any sentence (excluded before classification)", False, 12, MID_GRAY),
        ("Avg Conf = mean model confidence across all classified sentences for that protein", False, 12, MID_GRAY),
    ]
    add_multiline_box(slide, Inches(0.4), Inches(6.4), Inches(12), Inches(1), lines_excl)

    # =========================================================================
    # SLIDE 7: Cross-Domain Training
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Cross-Domain Training", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    lines_insight = [
        ("Key insight: association language is domain-general", True, 16, ACCENT_BLUE),
        ("", False, 6, DARK_TEXT),
        ("\"X is linked to Y\" and \"no correlation between X and Y\" -- these patterns are the same", False, 14, DARK_TEXT),
        ("regardless of whether X is a HFpEF protein or a cardiomyopathy protein.", False, 14, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Training exclusively on HFpEF sentences (n=992) limits the model unnecessarily.", False, 14, DARK_TEXT),
        ("A broader cardiology dataset should produce a larger, more balanced training set.", False, 14, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(1.2), Inches(11.5), Inches(2.5), lines_insight)

    cross_rows = [
        ["Experiment", "Training Data", "Strategy", "Rationale"],
        ["A: Cardio-only", "Broader cardiology\n(no HFpEF)", "Train from scratch", "Tests generalization: can the model learn\nassociation patterns without seeing HFpEF?"],
        ["B: Combined", "Cardiology + HFpEF", "Train from scratch", "Tests whether broader data + domain-specific\ndata outperforms either alone"],
        ["C: 2-Stage", "Cardiology then HFpEF", "Pre-train, fine-tune", "Broader patterns first, then domain-specific\nadjustment"],
    ]
    add_table(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(2.5), cross_rows,
              col_widths=[Inches(2.0), Inches(2.5), Inches(2.0), Inches(5.0)])

    add_text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
                 "All models evaluated on the same HFpEF held-out set (n=176) -- external validation",
                 font_size=14, color=MID_GRAY, bold=True)

    # =========================================================================
    # SLIDE 8: Results Table
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Results", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
                 "Held-out evaluation set: 176 samples (57 associated, 16 not_associated, 103 incidental)",
                 font_size=14, color=MID_GRAY)

    result_rows = [
        ["Model", "Accuracy", "Macro F1", "Weighted F1", "Notes"],
        ["PubMedBERT + Focal Loss", "69.3%", "0.656", "0.695", "Best model -- lr=1e-5, HFpEF augmented data"],
        ["Cardio + HFpEF Combined", "67.6%", "0.653", "0.673", "Broader cardiology + HFpEF data merged"],
        ["Ensemble (3 models)", "65.3%", "0.627", "0.655", "Softmax avg of top 3; weaker models diluted"],
        ["PubMedBERT (baseline)", "64.8%", "0.615", "0.647", "Default settings, no augmentation"],
        ["Cardio 2-Stage", "63.1%", "0.616", "0.633", "Broader pre-train -> HFpEF fine-tune"],
        ["Cardio-Only", "63.1%", "0.592", "0.635", "External validation: no HFpEF training data"],
        ["PubMedBERT 2-Stage (RE)", "61.4%", "0.585", "0.612", "ChemProt/DDI pre-train; domain shift too large"],
        ["Rule-based (spaCy)", "34.7%", "0.291", "0.295", "Pattern matching; confirms need for ML"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.8), result_rows,
              col_widths=[Inches(2.8), Inches(1.2), Inches(1.2), Inches(1.4), Inches(4.9)])

    # What worked / didn't
    lines_worked = [
        ("What worked", True, 16, GREEN_ACCENT),
        ("Focal loss: +4.5pp accuracy by focusing on hard/rare classes", False, 13, DARK_TEXT),
        ("Data augmentation: Expanding not_associated 89 -> 270 improved minority F1", False, 13, DARK_TEXT),
        ("Lower learning rate (1e-5): More stable fine-tuning for small datasets", False, 13, DARK_TEXT),
        ("Cross-domain combined (67.6%): Broader cardiology data closes gap with best", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(5.7), Inches(5.2), Inches(2), lines_worked)

    lines_didnt = [
        ("What didn't work", True, 16, RED_ACCENT),
        ("Cardio-only (63.1%): Domain transfer alone underperforms HFpEF-specific data", False, 13, DARK_TEXT),
        ("Cardio 2-stage (63.1%): No benefit over cardio-only; HFpEF fine-tune too brief", False, 13, DARK_TEXT),
        ("Ensemble: Averaging in weaker models diluted the best model", False, 13, DARK_TEXT),
        ("ChemProt/DDI pre-train: Domain shift too large for HFpEF", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(6.5), Inches(5.7), Inches(5.5), Inches(2), lines_didnt)

    # =========================================================================
    # SLIDE 9: Dashboard Figure
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Model Comparison Dashboard", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    dashboard_path = ROOT / "results" / "hfpef_v3_final_comparison.png"
    if dashboard_path.exists():
        slide.shapes.add_picture(str(dashboard_path),
                                 Inches(0.5), Inches(1.2), Inches(12.3), Inches(6.0))

    # =========================================================================
    # SLIDE 10: CaseOLAP Impact
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Impact on CaseOLAP Rankings", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # Before/after boxes
    # Before
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.5))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xFD, 0xED, 0xED)
    box.line.color.rgb = RED_ACCENT

    lines_before = [
        ("Before NLP Filtering", True, 18, RED_ACCENT),
        ("", False, 6, DARK_TEXT),
        ("3,554 proteins ranked by text co-occurrence", False, 14, DARK_TEXT),
        ("Top-ranked proteins included false positives:", False, 14, DARK_TEXT),
        ("  SPARC (score 0.847), Stanniocalcin-1 (0.847), Albumin (0.566)", False, 12, MID_GRAY),
        ("No distinction between meaningful vs. incidental mentions", False, 14, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(1.0), Inches(1.6), Inches(5.0), Inches(2.3), lines_before)

    # After
    box2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(7.0), Inches(1.5), Inches(5.5), Inches(2.5))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(0xE8, 0xF8, 0xF0)
    box2.line.color.rgb = GREEN_ACCENT

    lines_after = [
        ("After NLP Filtering", True, 18, GREEN_ACCENT),
        ("", False, 6, DARK_TEXT),
        ("500 proteins excluded (14% of total)", False, 14, DARK_TEXT),
        ("All had zero associated mentions in evidence sentences", False, 14, DARK_TEXT),
        ("High-confidence false positives caught and removed", False, 14, DARK_TEXT),
        ("Remaining proteins have sentence-level evidence", False, 14, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(7.2), Inches(1.6), Inches(5.0), Inches(2.3), lines_after)

    # Arrow between
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(6.35), Inches(2.5), Inches(0.6), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT_BLUE
    arrow.line.fill.background()

    # Filter module description
    lines_filter = [
        ("CaseOLAP Filter Module (src/caseolap_filter.py)", True, 15, ACCENT_BLUE),
        ("", False, 6, DARK_TEXT),
        ("Integrates directly into the CaseOLAP pipeline as a post-processing step.", False, 14, DARK_TEXT),
        ("For each protein: classifies all evidence sentences, tallies associated vs. not_associated vs. incidental.", False, 14, DARK_TEXT),
        ("Inclusion criterion: associated_mentions > (not_associated_mentions + incidental_mentions)", False, 14, DARK_TEXT),
        ("Returns per-protein confidence scores and specific exclusion reasons.", False, 14, DARK_TEXT),
        ("Processes all 3,554 proteins in minutes (vs. months of manual review).", False, 14, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(2.5), lines_filter)

    # =========================================================================
    # SLIDE 11: Next Steps
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                 "Next Steps", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    next_items = [
        ("1", "More Labeled Data", "The 992-sentence training set is the binding constraint. Another 300 labeled sentences, especially negative findings, would yield more improvement than any model architecture change.", "HIGH"),
        ("2", "GPU Training", "BiomedBERT-Large (340M parameters) could not be trained on laptop hardware. A server with a GPU would enable this and likely add 3-5 percentage points.", "MEDIUM"),
        ("3", "Active Learning", "Use model uncertainty to prioritize which sentences to label next. Focus expert labeling effort on cases where the model is least confident.", "MEDIUM"),
        ("4", "Human-in-the-Loop", "Deploy classifier on new CaseOLAP outputs, flag low-confidence predictions for expert review, continuously improve with feedback.", "FUTURE"),
    ]

    y = Inches(1.5)
    for num, title, desc, priority in next_items:
        # Priority badge
        p_color = GREEN_ACCENT if priority == "HIGH" else (ORANGE_ACCENT if priority == "MEDIUM" else MID_GRAY)
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(11.0), y + Inches(0.1), Inches(1.2), Inches(0.35))
        badge.fill.solid()
        badge.fill.fore_color.rgb = p_color
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].text = priority
        tf.paragraphs[0].font.size = Pt(10)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Number
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                        Inches(0.8), y, Inches(0.55), Inches(0.55))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        add_text_box(slide, Inches(1.6), y - Inches(0.05), Inches(9), Inches(0.4),
                     title, font_size=17, color=DARK_TEXT, bold=True)
        add_text_box(slide, Inches(1.6), y + Inches(0.4), Inches(9.2), Inches(0.8),
                     desc, font_size=13, color=MID_GRAY)
        y += Inches(1.4)

    # =========================================================================
    # SLIDE 12: Separator (Old vs New)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(1), Inches(2.1), Inches(11.5), Inches(1.0),
                 "NEW SECTION: Requested Additions",
                 font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

    add_text_box(slide, Inches(1), Inches(3.3), Inches(11.5), Inches(0.9),
                 "Slides below document what was newly implemented on the server run",
                 font_size=20, color=ACCENT_TEAL, bold=False, alignment=PP_ALIGN.LEFT)

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(3.05), Inches(4.8), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(1), Inches(5.1), Inches(11), Inches(0.5),
                 "Legacy content above this slide | New updates below",
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

    # =========================================================================
    # SLIDE 13: Update 1 (Disease Expansion + Mining)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11.8), Inches(0.7),
                 "Update 1: Disease Expansion Implemented", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.25), Inches(12.0), Inches(0.5),
                 "Requested set implemented: CVA, IHD, CM, CHD, ARR, VD plus HFpEF",
                 font_size=14, color=MID_GRAY)

    disease_rows = [
        ["Code", "Disease", "Mined Co-mention Sentences"],
        ["CVA", "Stroke / Cerebrovascular Disease", "5,212"],
        ["IHD", "Ischemic Heart Disease", "3,350"],
        ["CM", "Cardiomyopathies", "1,983"],
        ["CHD", "Congenital Heart Disease", "1,555"],
        ["ARR", "Cardiac Arrhythmias", "2,101"],
        ["VD", "Valve Disease", "2,486"],
        ["Total", "Cross-disease mined corpus", "16,687"],
    ]
    add_table(slide, Inches(0.8), Inches(1.9), Inches(11.6), Inches(3.8), disease_rows,
              col_widths=[Inches(1.2), Inches(5.4), Inches(3.0)])

    lines_cvd_status = [
        ("Server outputs (Feb 14, 2026): cvd_broad_train.json = 3,446 | cvd_hfpef_combined_train.json = 4,619", False, 13, DARK_TEXT),
        ("CVD labels: 379 associated | 67 not_associated | 3000 incidental", False, 13, DARK_TEXT),
        ("Combined labels: 701 associated | 337 not_associated | 3581 incidental", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(6.0), Inches(12.0), Inches(1.2), lines_cvd_status)

    # =========================================================================
    # SLIDE 14: Method 1 (Sentence Baseline)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Method 1: Sentence-Only Baseline", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.0), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    lines_m1 = [
        ("Configuration", True, 16, ACCENT_BLUE),
        ("PubMedBERT + focal loss on hfpef_v3_train_augmented.json (1,173 samples)", False, 13, DARK_TEXT),
        ("Input uses only the target protein-disease sentence", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(1.35), Inches(11.8), Inches(1.4), lines_m1)

    method1_rows = [
        ["Metric", "Value", "Comment"],
        ["Accuracy", "69.3%", "Strong baseline for sentence-only evidence"],
        ["Macro F1", "0.6667", "Minority classes remain difficult"],
        ["Weighted F1", "0.6929", "Best single-model sentence performance"],
    ]
    add_table(slide, Inches(0.8), Inches(3.0), Inches(11.6), Inches(2.2), method1_rows,
              col_widths=[Inches(2.0), Inches(1.7), Inches(7.5)])

    add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.9), Inches(0.9),
                 "Method 1 is the reference point for context-window improvements.",
                 font_size=13, color=MID_GRAY)

    # =========================================================================
    # SLIDE 15: Method 2 (Context Window Model)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Method 2: Context-Window Model", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.0), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    pipeline_rows = [
        ["Step", "Implementation Detail", "Output"],
        ["1", "Extract target sentence with protein + disease co-mention", "Sentence sample"],
        ["2", "Attach context window (target +/- 1 sentence)", "Context-enriched sample"],
        ["3", "Apply fuzzy-matching fallback for context alignment", "Stable mapping"],
        ["4", "Train dedicated context model (pubmedbert_context)", "Context-specific encoder"],
    ]
    add_table(slide, Inches(0.8), Inches(1.7), Inches(11.6), Inches(2.9), pipeline_rows,
              col_widths=[Inches(0.8), Inches(6.8), Inches(4.0)])

    method2_rows = [
        ["Evaluation", "Accuracy", "Macro F1", "Weighted F1"],
        ["Context model on context eval", "66.5%", "0.6425", "0.6643"],
        ["Context model on sentence eval", "56.8%", "0.4480", "0.5300"],
        ["Sentence model on context eval", "60.8%", "--", "--"],
    ]
    add_table(slide, Inches(0.8), Inches(4.9), Inches(11.6), Inches(2.3), method2_rows,
              col_widths=[Inches(4.3), Inches(1.5), Inches(1.5), Inches(1.8)])

    # =========================================================================
    # SLIDE 16: Method 3 (Fusion + Threshold Tuning)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Method 3: Dual-View Fusion + Threshold Tuning", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.6), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    lines_m3 = [
        ("Fusion rule: P = 0.7 * sentence_probs + 0.3 * context_probs", True, 15, ACCENT_BLUE),
        ("Threshold tuning reduces incidental overprediction", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(1.25), Inches(11.8), Inches(1.1), lines_m3)

    method3_rows = [
        ["Variant", "Accuracy", "Macro F1", "Weighted F1", "Notes"],
        ["Fusion (argmax)", "70.5%", "0.6793", "0.7034", "alpha_context=0.3"],
        ["Fusion + threshold tuning", "71.0%", "0.6909", "0.7095", "Best benchmark from Feb 14 server run"],
    ]
    add_table(slide, Inches(0.8), Inches(2.7), Inches(11.6), Inches(1.9), method3_rows,
              col_widths=[Inches(3.4), Inches(1.4), Inches(1.4), Inches(1.6), Inches(3.8)])

    threshold_rows = [
        ["Class", "Best Threshold"],
        ["associated", "0.20"],
        ["not_associated", "0.45"],
        ["incidental", "0.30"],
    ]
    add_table(slide, Inches(0.8), Inches(4.9), Inches(4.4), Inches(1.9), threshold_rows,
              col_widths=[Inches(2.5), Inches(1.5)])

    add_text_box(slide, Inches(5.5), Inches(5.0), Inches(6.8), Inches(1.4),
                 "Current best option for deployment: Method 3 with thresholds.\nImprovement vs Method 1: +1.7 points accuracy.",
                 font_size=13, color=DARK_TEXT)

    # =========================================================================
    # SLIDE 17: Dataset and Sentence Growth
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Update: Sentence and Dataset Growth", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    growth_rows = [
        ["Stage", "Train Samples", "Delta vs 992 Baseline", "What Changed"],
        ["HFpEF v3 baseline", "992", "Baseline", "Manually labeled sentence-level set"],
        ["HFpEF + hard negatives", "1,173", "+181 (+18.2%)", "not_associated expanded: 89 -> 270"],
        ["HFpEF + context view", "1,173", "Representation change", "Target sentence plus local context"],
        ["CVD broad (server)", "3,446", "+2,454 (+247.4%)", "Pseudo-labeled mining across six diseases"],
        ["CVD + HFpEF combined (server)", "4,619", "+3,627 (+365.6%)", "Merged broad CVD with HFpEF seed"],
    ]
    add_table(slide, Inches(0.8), Inches(1.7), Inches(11.8), Inches(3.8), growth_rows,
              col_widths=[Inches(3.2), Inches(1.5), Inches(2.4), Inches(4.4)])

    add_text_box(slide, Inches(0.8), Inches(5.8), Inches(11.8), Inches(0.9),
                 "Scale-up succeeded; next gains depend on context selection quality and calibration.",
                 font_size=13, color=MID_GRAY)

    # =========================================================================
    # SLIDE 18: Confidence-Based Filtering Output
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Update: Confidence Reporting for Exclusion", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    conf_rows = [
        ["Requested Item", "Status", "Details"],
        ["Confidence-focused exclusion output", "Implemented", "Per-sentence confidence and avg_confidence are available."],
        ["Current excluded count", "500 proteins (prior snapshot)", "500 / 3,554 proteins excluded in existing report export."],
        ["Next recompute step", "Pending rerun", "Rerun CaseOLAP filter using Method 3 fused model."],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(2.7), conf_rows,
              col_widths=[Inches(3.5), Inches(2.3), Inches(6.0)])

    add_text_box(slide, Inches(0.8), Inches(4.9), Inches(11.8), Inches(1.4),
                 "PI-facing export columns: protein_id, associated_count, not_associated_count, incidental_count, association_confidence, avg_confidence.",
                 font_size=13, color=DARK_TEXT)

    # =========================================================================
    # SLIDE 19: Methods 4-7 Results (Implemented)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Research-Inspired Methods: Measured Results", font_size=28, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(4.2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(1.25), Inches(12.0), Inches(0.5),
                 "Metrics source: results/context_window_research_methods_20260215_220500.json",
                 font_size=12, color=MID_GRAY)

    method_rows = [
        ["Method", "Implementation", "Accuracy", "Macro F1", "Delta vs Method 1 (Macro F1)"],
        ["Method 1", "Sentence-only baseline", "69.3%", "0.6667", "baseline"],
        ["Method 4", "Position-aware context ordering", "68.2%", "0.6696", "+0.0029"],
        ["Method 5", "Entity-aware top-k context selection", "68.8%", "0.6683", "+0.0017"],
        ["Method 6", "Two-pass reranking on uncertain samples", "68.8%", "0.6762", "+0.0096"],
        ["Method 7", "Adaptive context budget by uncertainty", "66.5%", "0.6534", "-0.0133"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(3.5), method_rows,
              col_widths=[Inches(1.5), Inches(4.7), Inches(1.4), Inches(1.4), Inches(2.6)])

    lines_m4 = [
        ("Best new method by macro F1: Method 6 (two-pass rerank).", True, 14, ACCENT_BLUE),
        ("Interpretation: context helps hardest/uncertain cases more than easy cases.", False, 13, DARK_TEXT),
        ("Observed tradeoff: highest macro F1 is not the highest overall accuracy.", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(5.5), Inches(11.8), Inches(1.3), lines_m4)

    # =========================================================================
    # SLIDE 20: References Mapped to Implemented Methods
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Newest NLP References and What Was Implemented", font_size=28, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(4.4), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    ref_rows = [
        ["Reference", "Main Finding", "Implemented as"],
        ["Liu et al., TACL 2024", "Middle-position evidence is often underused", "Method 4: target-at-edges context ordering"],
        ["Zhang et al., NeurIPS 2024", "Position handling improves middle retrieval", "Method 4/5: reordered and scored context"],
        ["Yue et al., ICLR 2025", "Inference compute scaling helps long context", "Method 6: second pass only for uncertain samples"],
        ["Jin et al., ICLR 2025", "Extra retrieved evidence can add hard negatives", "Method 7: adaptive top-k context budget"],
        ["Li et al., Findings EMNLP 2025", "Entity type helps evidence retrieval", "Method 5: entity-aware sentence scoring"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(3.8), ref_rows,
              col_widths=[Inches(2.6), Inches(4.3), Inches(4.9)])

    add_text_box(slide, Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.9),
                 "All listed methods were implemented in scripts/evaluation/context_window_research_methods.py and benchmarked on hfpef_v3_eval (n=176).",
                 font_size=12, color=MID_GRAY)

    add_text_box(
        slide,
        Inches(0.8), Inches(6.65), Inches(12.0), Inches(0.5),
        "Refs: aclanthology.org/2024.tacl-1.9 | papers.nips.cc/.../6ffdbbe... | proceedings.iclr.cc/.../b574717... | proceedings.iclr.cc/.../5df5b1f... | aclanthology.org/2025.findings-emnlp.961/",
        font_size=9, color=MID_GRAY,
    )

    # =========================================================================
    # SLIDE 21: New Experiment Block Separator
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)

    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.5), Inches(1.0),
                 "NEW BLOCK: 80%-Target Experiments",
                 font_size=38, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, Inches(1), Inches(3.25), Inches(11.5), Inches(0.8),
                 "All slides below are newly added and do not modify earlier content",
                 font_size=20, color=ACCENT_TEAL, alignment=PP_ALIGN.LEFT)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(1), Inches(3.0), Inches(5.4), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    # =========================================================================
    # SLIDE 22: Data Pipeline Additions
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "New Data Pipeline Implementations", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.5), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    data_rows = [
        ["Component", "Script", "Output"],
        ["Disagreement hardcase mining", "scripts/data_prep/mine_disagreement_hardcases.py", "hardcases_train_disagreement_20260217_204501.json (600)"],
        ["not_associated expansion", "scripts/data_prep/build_not_associated_expansion.py", "hfpef_v5_train_not_assoc_expanded.json (1227); not_associated 270 -> 326"],
        ["Consensus pseudo-labeling", "scripts/data_prep/build_consensus_pseudolabels.py", "hfpef_v5_pseudo_consensus.json (0 selected at confidence >= 0.95)"],
        ["Training mix builder", "scripts/data_prep/build_training_mix.py", "hfpef_v5_train_mix.json (1225 after duplicate drop)"],
        ["Dedicated not_associated expert", "scripts/training/train_not_assoc_expert.py", "models/hfpef_v5/not_assoc_expert.joblib (val macro F1 0.9408)"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(3.8), data_rows,
              col_widths=[Inches(2.5), Inches(4.6), Inches(4.7)])

    add_text_box(slide, Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.9),
                 "v5 run focus: targeted class-balance correction and expert rerouting; broad pseudo-labeling gate accepted 0/300 candidates.",
                 font_size=13, color=DARK_TEXT)

    # =========================================================================
    # SLIDE 23: 80%-Target Experiment Results
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "80%-Target Experiments: Measured Outcomes", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.8), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    result_rows_new = [
        ["Method", "Accuracy", "Macro F1", "Status"],
        ["Baseline (Method 1 sentence)", "69.3%", "0.6667", "Reference baseline"],
        ["Fusion + threshold tuning (server)", "74.4%", "0.6789", "Previous best before v5"],
        ["v5 not_associated expert cascade", "75.0%", "0.6935", "Current best overall (assoc_max=0.35, not_inc_gap=0.03)"],
        ["v5 dual-expert cascade", "72.7%", "0.7111", "Best macro F1; accuracy tradeoff (rerouted_assoc=28)"],
        ["v5 mix retrain (expanded + pseudo mix)", "66.5%", "0.6654", "Regression vs fusion baseline"],
        ["NLI reformulation", "62.5%", "0.4686", "Large associated-class instability"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(3.5), result_rows_new,
              col_widths=[Inches(3.6), Inches(1.3), Inches(1.3), Inches(5.1)])

    add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11.8), Inches(1.0),
                 "Key finding: best accuracy is 75.0%, while best macro F1 is 0.7111 from dual-expert rerouting; 80% has not been reached yet.",
                 font_size=13, color=RED_ACCENT, bold=True)

    # =========================================================================
    # SLIDE 24: Bottleneck and Next Run Plan
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Why 80% Is Not Reached Yet", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    bottleneck_rows = [
        ["Bottleneck", "Observed Evidence", "Next Fix"],
        ["Accuracy vs macro-F1 tradeoff", "Dual-expert run reached macro F1 0.7111 but dropped accuracy to 72.7%", "Tune reroute budget to cap incidental false positives while keeping associated gains"],
        ["Pseudo-label gate too strict", "Consensus pseudo-labeling accepted 0/300 candidates", "Use disease-conditioned prompts or lower threshold only after manual spot-check"],
        ["Noise sensitivity in retraining", "v5 mixed retrain fell to 66.5% / 0.6654", "Filter added samples by sentence quality and agreement before retraining"],
        ["Expert reroute impact is narrow", "Best gain came from rerouting just 1 sample", "Train multiclass expert or larger uncertainty band for rerouting"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(3.8), bottleneck_rows,
              col_widths=[Inches(2.8), Inches(3.9), Inches(5.1)])

    add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.8), Inches(0.8),
                 "Most defensible near-term target: 76-78% with associated-focused mining and stricter sample-quality controls.",
                 font_size=13, color=ACCENT_BLUE, bold=True)

    # =========================================================================
    # SLIDE 25: AWS v8 Run Summary (New)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "NEW: AWS v8 Run Summary (Feb 23, 2026)", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(4.3), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    lines_v8 = [
        ("Executed pipeline on AWS server: ec2-18-224-98-201.us-east-2.compute.amazonaws.com", False, 13, DARK_TEXT),
        ("Objective: rerun full \"do all\" block from pseudo build through no-leak evaluation.", False, 13, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Completed stages", True, 15, ACCENT_BLUE),
        ("1) CVD full-UniProt quality pseudo pool build", False, 13, DARK_TEXT),
        ("2) Precision-controlled pseudo pass over hfpef_corpus", False, 13, DARK_TEXT),
        ("3) Training mix + context-mapped mix generation", False, 13, DARK_TEXT),
        ("4) Three model retrains (sentence/context/cvd)", False, 13, DARK_TEXT),
        ("5) Expert retrains (not_associated + associated)", False, 13, DARK_TEXT),
        ("6) No-leak eval: calibration + not-assoc cascade + dual cascade", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(0.8), Inches(1.4), Inches(11.8), Inches(4.4), lines_v8)

    add_text_box(slide, Inches(0.8), Inches(6.1), Inches(11.8), Inches(0.8),
                 "Artifacts synced locally: logs/aws_v8_20260223/*",
                 font_size=12, color=MID_GRAY)

    # =========================================================================
    # SLIDE 26: AWS v8 Data Volumes (New)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "NEW: AWS v8 Data Volumes", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(3.2), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    v8_data_rows = [
        ["Source / Stage", "Count", "Class Distribution / Notes"],
        ["CVD full-UniProt quality pool", "741", "assoc 205 | not_assoc 36 | incidental 500"],
        ["Precision pseudo pool", "0", "strict gate accepted 0 / 15,989 candidates"],
        ["Pseudo combined", "741", "same as quality pool (no extra precision pseudo accepted)"],
        ["Base relabel2 train", "1,227", "assoc 371 | not_assoc 308 | incidental 548"],
        ["v8 train mix output", "1,587", "assoc 534 | not_assoc 342 | incidental 711"],
        ["Context-mapped mix", "1,587", "mapped 1,060 | fallback 527"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(3.9), v8_data_rows,
              col_widths=[Inches(3.7), Inches(1.4), Inches(6.7)])

    add_text_box(slide, Inches(0.8), Inches(5.95), Inches(11.8), Inches(0.8),
                 "Key blocker from this run: precision pseudo stage added zero new samples.",
                 font_size=13, color=RED_ACCENT, bold=True)

    # =========================================================================
    # SLIDE 27: AWS v8 No-Leak Results (New)
    # =========================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)

    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "NEW: AWS v8 No-Leak Results vs Relabel2", font_size=30, color=DARK_TEXT, bold=True)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0.8), Inches(1.0), Inches(4.0), Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()

    v8_metric_rows = [
        ["Variant", "Accuracy", "Macro F1", "Weighted F1"],
        ["Relabel2 calibrated (best prior)", "83.22%", "0.7385", "0.8317"],
        ["v8 calibrated (best v8)", "82.52%", "0.7427", "0.8298"],
        ["v8 not-assoc cascade", "79.72%", "0.7362", "0.8129"],
        ["v8 dual cascade", "68.53%", "0.6610", "0.7156"],
    ]
    add_table(slide, Inches(0.8), Inches(1.8), Inches(11.8), Inches(2.9), v8_metric_rows,
              col_widths=[Inches(4.2), Inches(1.6), Inches(1.6), Inches(1.9)])

    delta_rows = [
        ["Best-v8 minus best-relabel2", "Delta"],
        ["Accuracy delta", "-0.699 pp"],
        ["Macro F1 delta", "+0.421 pp"],
        ["Weighted F1 delta", "-0.198 pp"],
    ]
    add_table(slide, Inches(0.8), Inches(5.0), Inches(4.8), Inches(1.8), delta_rows,
              col_widths=[Inches(3.3), Inches(1.3)])

    lines_v8_takeaway = [
        ("Takeaway", True, 14, ACCENT_BLUE),
        ("v8 slightly improves macro F1 but does not improve top-line accuracy.", False, 13, DARK_TEXT),
        ("Next accuracy gain likely needs accepted high-quality pseudo samples (not zero-pass gating).", False, 13, DARK_TEXT),
    ]
    add_multiline_box(slide, Inches(5.9), Inches(5.0), Inches(6.7), Inches(1.8), lines_v8_takeaway)

    # Save
    output_path = ROOT / "results" / "hfpef_presentation.pptx"
    prs.save(str(output_path))
    print(f"Saved presentation to {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
