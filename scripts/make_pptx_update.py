"""Generate focused HFpEF update presentation for Clodomir meeting."""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent.parent

DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x7A, 0xCC)
ACCENT_TEAL = RGBColor(0x00, 0xB4, 0xD8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MID_GRAY = RGBColor(0xA0, 0xA0, 0xA0)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
TABLE_HEADER_BG = RGBColor(0x00, 0x56, 0x8A)
TABLE_ROW_LIGHT = RGBColor(0xF5, 0xF9, 0xFC)
TABLE_ROW_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED_ACCENT = RGBColor(0xE7, 0x4C, 0x3C)
GREEN_ACCENT = RGBColor(0x2E, 0xCC, 0x71)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    return txBox


def add_multiline_box(slide, left, top, width, height, lines, font_size=13,
                      color=DARK_TEXT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, bold, sz, clr) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz if sz else font_size)
        p.font.color.rgb = clr if clr else color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = Pt(4)
    return txBox


def add_table(slide, left, top, width, height, rows, col_widths=None):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), left, top, width, height)
    table = table_shape.table
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                if row_idx == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                    paragraph.alignment = PP_ALIGN.CENTER
                else:
                    paragraph.font.color.rgb = DARK_TEXT
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_WHITE
    return table_shape


def accent_bar(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # =================================================================
    # SLIDE 1: Title
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.2),
                 "HFpEF Protein-Disease\nAssociation Classifier",
                 font_size=42, color=WHITE, bold=True)
    accent_bar(slide, Inches(1), Inches(3.7), Inches(4))
    add_text_box(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
                 "Progress Update -- April 2026",
                 font_size=22, color=ACCENT_TEAL)
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
                 "Aktan Azat  |  Division of Cardiovascular Medicine, UC Davis",
                 font_size=16, color=MID_GRAY)

    # =================================================================
    # SLIDE 2: Recap -- The Problem
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Recap: The Problem", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(2.5))

    add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.0), [
        ("CaseOLAP ranks proteins by co-occurrence with HFpEF in PubMed.", False, 16, DARK_TEXT),
        ("", False, 8, DARK_TEXT),
        ("Co-occurrence is not association.", True, 18, RED_ACCENT),
        ("", False, 8, DARK_TEXT),
        ("A protein mentioned in patient demographics, exclusion criteria, or background", False, 14, DARK_TEXT),
        ("gets the same score as one the paper actually studies.", False, 14, DARK_TEXT),
        ("", False, 8, DARK_TEXT),
        ("We need a filter that reads each sentence and decides:", False, 14, DARK_TEXT),
        ("is the paper claiming a real protein-disease relationship?", True, 15, ACCENT_BLUE),
    ])

    # =================================================================
    # SLIDE 3: Recap -- The Pipeline
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Recap: The Pipeline", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(2.5))

    steps = [
        ("1", "Define HFpEF", "MeSH terms for heart failure with preserved ejection fraction"),
        ("2", "Search PubMed", "Retrieve abstracts mentioning HFpEF and protein terms"),
        ("3", "Extract sentences", "spaCy splits abstracts; keep sentences with protein mentions"),
        ("4", "Classify", "NLP model labels each sentence: associated / not associated / incidental"),
        ("5", "Re-rank proteins", "Proteins scored by association evidence, not just co-occurrence"),
    ]
    step_rows = [["Step", "Stage", "What happens"]] + [[s[0], s[1], s[2]] for s in steps]
    add_table(slide, Inches(0.8), Inches(1.6), Inches(11.8), Inches(3.0), step_rows,
              col_widths=[Inches(0.8), Inches(2.5), Inches(8.5)])

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.8), Inches(0.6),
                 "Step 4 is the new piece. Everything else was already in place from CaseOLAP.",
                 font_size=13, color=MID_GRAY)

    # =================================================================
    # SLIDE 4: Recap -- Data and Model
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Recap: Data and Model", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(2.5))

    add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4.5), [
        ("Training data", True, 16, ACCENT_BLUE),
        ("~1,100 manually labeled sentences from PubMed", False, 14, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Three classes:", True, 14, DARK_TEXT),
        ("Associated (32%) -- paper claims a relationship", False, 13, DARK_TEXT),
        ("Not associated (9%) -- paper explicitly denies a link", False, 13, DARK_TEXT),
        ("Incidental (58%) -- protein just mentioned in passing", False, 13, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Class imbalance handled with focal loss and class weighting.", False, 13, MID_GRAY),
    ])

    add_multiline_box(slide, Inches(6.8), Inches(1.5), Inches(5.7), Inches(4.5), [
        ("Model", True, 16, ACCENT_BLUE),
        ("PubMedBERT -- pre-trained on PubMed abstracts", False, 14, DARK_TEXT),
        ("Fine-tuned on our labeled sentences", False, 14, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Training enhancements:", True, 14, DARK_TEXT),
        ("Focal loss -- pays more attention to hard examples", False, 13, DARK_TEXT),
        ("R-Drop -- regularization for stability across seeds", False, 13, DARK_TEXT),
        ("Calibrated thresholds -- tuned per class", False, 13, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("~20 minutes to train on a GPU.", False, 13, MID_GRAY),
    ])

    # =================================================================
    # SLIDE 5: Where We Were
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Where We Were Before", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(2.8))

    add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(2.5), [
        ("We tried many approaches to improve accuracy:", False, 15, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Context windows -- show the model surrounding sentences from the abstract", False, 14, DARK_TEXT),
        ("Fusion -- combine sentence-only and context models", False, 14, DARK_TEXT),
        ("Expert cascades -- specialized models for hard classes", False, 14, DARK_TEXT),
        ("Pseudo-labeling -- automatically label more training data from related papers", False, 14, DARK_TEXT),
        ("NLI reformulation -- reframe classification as a yes/no question", False, 14, DARK_TEXT),
    ])

    result_rows = [
        ["Approach", "Accuracy", "Macro F1"],
        ["Sentence-only baseline", "69.3%", "66.7%"],
        ["Fusion + threshold tuning", "74.4%", "67.9%"],
        ["Expert cascade (best)", "75.0%", "69.4%"],
        ["Calibrated (best overall)", "83.2%", "73.9%"],
    ]
    add_table(slide, Inches(0.8), Inches(4.3), Inches(6.5), Inches(2.4), result_rows,
              col_widths=[Inches(3.5), Inches(1.5), Inches(1.5)])

    add_text_box(slide, Inches(7.8), Inches(4.8), Inches(4.8), Inches(1.0),
                 "Best accuracy: 83.2%\nTarget: 80% -- met, but we wanted more.\nNo architecture change could push past ~83%.",
                 font_size=14, color=RED_ACCENT, bold=True)

    # =================================================================
    # SLIDE 6: Separator -- What Changed
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK_BG)
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.5), Inches(1.0),
                 "What Changed",
                 font_size=42, color=WHITE, bold=True)
    accent_bar(slide, Inches(1), Inches(3.0), Inches(3.5))
    add_multiline_box(slide, Inches(1), Inches(3.4), Inches(11.5), Inches(2.5), [
        ("We stopped trying to fix the models.", False, 20, ACCENT_TEAL),
        ("We looked at the labels instead.", False, 20, ACCENT_TEAL),
    ])

    # =================================================================
    # SLIDE 7: Label Audit -- Round 1
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Label Audit: Round 1 -- Evaluation Set", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3.8))

    add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.0), [
        ("All 6 trained models disagreed with 88 of 300 eval labels (29%).", False, 15, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("We sent these 88 sentences to Clodomir for expert review.", False, 15, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Problem discovered:", True, 15, ACCENT_BLUE),
        ("Many sentences mention multiple proteins. Without specifying which protein", False, 14, DARK_TEXT),
        ("we are asking about, even an expert can label the same sentence differently.", False, 14, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Fix: added a protein name column to the review spreadsheet.", False, 14, DARK_TEXT),
        ("This resolved most of the ambiguity.", False, 14, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Result: 53 of 88 labels corrected.", True, 15, DARK_TEXT),
        ("Models immediately jumped from ~70% to 85.7% accuracy -- with no retraining.", True, 15, GREEN_ACCENT),
    ])

    # =================================================================
    # SLIDE 8: Label Audit -- Round 2
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Label Audit: Round 2 -- Training Set", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3.8))

    add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(3.0), [
        ("164 suspicious training examples flagged.", False, 15, DARK_TEXT),
        ("Clodomir reviewed and corrected 109 labels.", False, 15, DARK_TEXT),
        ("", False, 6, DARK_TEXT),
        ("Models retrained on corrected training set", False, 15, DARK_TEXT),
        ("(1,102 examples, 3 seeds per configuration).", False, 15, DARK_TEXT),
    ])

    transition_rows = [
        ["Label Change", "Count"],
        ["incidental -> associated", "57"],
        ["associated -> incidental", "27"],
        ["associated -> not_associated", "10"],
        ["incidental -> not_associated", "9"],
        ["not_associated -> incidental", "4"],
        ["not_associated -> associated", "2"],
    ]
    add_table(slide, Inches(7.2), Inches(1.5), Inches(5.3), Inches(2.6), transition_rows,
              col_widths=[Inches(3.5), Inches(1.3)])

    add_multiline_box(slide, Inches(0.8), Inches(4.8), Inches(11.8), Inches(2.0), [
        ("Dominant correction: incidental -> associated (57 of 109).", True, 14, ACCENT_BLUE),
        ("Many real protein-disease associations had been labeled as incidental mentions.", False, 14, DARK_TEXT),
        ("This systematic bias was suppressing model performance on the associated class.", False, 14, DARK_TEXT),
    ])

    # =================================================================
    # SLIDE 9: Results After Corrections
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Results After Label Corrections", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3.5))

    result_rows = [
        ["Model", "Accuracy", "Macro F1"],
        ["Best calibrated single model", "86.3%", "87.1%"],
        ["6-model ensemble", "86.3%", "87.0%"],
        ["3-seed R-Drop ensemble", "86.0%", "86.7%"],
        ["", "", ""],
        ["Best before correction", "83.2%", "73.9%"],
    ]
    add_table(slide, Inches(0.8), Inches(1.6), Inches(7.0), Inches(2.8), result_rows,
              col_widths=[Inches(3.5), Inches(1.5), Inches(1.5)])

    add_multiline_box(slide, Inches(8.3), Inches(1.6), Inches(4.5), Inches(2.8), [
        ("Accuracy: 83.2% -> 86.3%", True, 18, GREEN_ACCENT),
        ("", False, 8, DARK_TEXT),
        ("Macro F1: 73.9% -> 87.1%", True, 18, GREEN_ACCENT),
        ("", False, 8, DARK_TEXT),
        ("Same model architecture.", False, 14, DARK_TEXT),
        ("Only the labels changed.", False, 14, DARK_TEXT),
    ])

    add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11.8), Inches(1.2),
                 "The macro F1 jump (+13 points) is especially important -- it means the model improved substantially on the minority classes, not just the majority class.",
                 font_size=14, color=DARK_TEXT)

    # =================================================================
    # SLIDE 10: Key Finding
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Key Finding", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(2.0))

    add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.5), [
        ("The models were already performing well.", True, 18, ACCENT_BLUE),
        ("They were being scored against labels that were, in many cases, wrong.", False, 15, DARK_TEXT),
        ("", False, 8, DARK_TEXT),
        ("Root cause of label noise:", True, 15, DARK_TEXT),
        ("Sentences with multiple proteins created ambiguity about which protein-disease", False, 14, DARK_TEXT),
        ("pair was being classified. Without a target protein column, labelers had to guess.", False, 14, DARK_TEXT),
        ("", False, 8, DARK_TEXT),
        ("Clodomir's feedback:", True, 15, DARK_TEXT),
        ("\"Labelling those sentences was not easy because there's a lot of ambiguity", False, 14, MID_GRAY),
        ("and subjectivity involved.\"", False, 14, MID_GRAY),
        ("", False, 8, DARK_TEXT),
        ("What this means:", True, 15, DARK_TEXT),
        ("For specialized biomedical classification, label quality matters more than", False, 14, DARK_TEXT),
        ("model architecture. We spent months improving models when the data needed fixing.", False, 14, DARK_TEXT),
    ])

    # =================================================================
    # SLIDE 11: Remaining Errors and Next Steps
    # =================================================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.7),
                 "Remaining Errors and Next Steps", font_size=30, color=DARK_TEXT, bold=True)
    accent_bar(slide, Inches(0.8), Inches(1.0), Inches(3.3))

    add_multiline_box(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(5.0), [
        ("Remaining ~14% of errors", True, 16, ACCENT_BLUE),
        ("Concentrated at the associated / incidental boundary.", False, 14, DARK_TEXT),
        ("Cases where the sentence is ambiguous even with the target protein specified.", False, 14, DARK_TEXT),
        ("Reasonable experts could disagree on these.", False, 14, DARK_TEXT),
        ("", False, 10, DARK_TEXT),
        ("To go further", True, 16, ACCENT_BLUE),
        ("Tighter annotation guidelines for the associated-incidental boundary.", False, 14, DARK_TEXT),
        ("Define what counts as 'association' more precisely -- does a correlational finding", False, 14, DARK_TEXT),
        ("count, or only mechanistic evidence?", False, 14, DARK_TEXT),
        ("", False, 10, DARK_TEXT),
        ("Apply the classifier to the full CaseOLAP output and generate updated", False, 14, DARK_TEXT),
        ("protein rankings for HFpEF.", False, 14, DARK_TEXT),
        ("", False, 10, DARK_TEXT),
        ("That last question is for the domain experts -- it's not a modeling decision.", False, 14, MID_GRAY),
    ])

    # Save
    output_path = ROOT / "results" / "hfpef_update_presentation.pptx"
    prs.save(str(output_path))
    print(f"Saved to {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
